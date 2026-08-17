"""SUPERSEDED / FALLBACK ONLY — the live master template is no longer this
script's output. The pipeline now duplicates a copy of the real, previously-
used Interactive Keynote deck ("INTERACTIVE KEYNOTE (CAS X PELLO)",
originally at drive file 1BnT1b5G5m1gOXI7UcJrbUdfaQvCNlC-Kcp9BAkp0Jt8, copied
into the Templates Drive folder as "Interactive Keynote Template - Crossing
the Ice" with its 5 client-logo image instances tagged alt-text
"client_logo") — which has the real design, photography, and layout this
script could only approximate. See AS_BUILT_interactive_keynote_proposal_
automation.txt for how that real template was identified and tagged.

This script is kept only as a documented fallback generator in case the real
Drive template file is ever lost/unavailable — it builds a best-effort .pptx
recreation from INTERACTIVE_KEYNOTE_reference_proposal.pdf's extracted text,
using bracket tokens ([CLIENT_ORG], [EVENT_DATE], [VENUE], [AUDIENCE_SIZE],
[DURATION]) for dynamic spots. slides_rewriter.py's current prompt targets
the real deck's actual personalization surface (one client-name mention in
the bio text, plus 5 tagged logo images) and does NOT look for these bracket
tokens — if this fallback template is ever put into production, the rewrite
prompt in src/slides_rewriter.py would need to be reverted to the
bracket-token-driven version to match it.

No PDF renderer was available in this environment, so this recreates
structure/wording/section order faithfully but cannot reproduce the
reference deck's exact background imagery, fonts, or color values — treat
the visual design below as a best-effort placeholder, not a pixel-perfect
match.

Usage:
    python build_master_template.py                # writes the .pptx only
    python build_master_template.py --upload        # also uploads to Drive
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BASE_DIR = Path(__file__).resolve().parent
OUTPUT_PATH = BASE_DIR / "templates" / "interactive_keynote_master_template.pptx"
LOGO_PLACEHOLDER_PATH = BASE_DIR / "templates" / "_logo_placeholder.png"

# 16:9 widescreen, matching a modern Keynote/PowerPoint export
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x1F, 0x33)
ICE_BLUE = RGBColor(0x7E, 0xC8, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x55, 0x55, 0x55)

FONT = "Montserrat"
BODY_FONT = "Calibri"


def _make_logo_placeholder_image() -> Path:
    """Generates a simple placeholder graphic for the client-logo image
    shape, since no real client logo exists until a proposal is generated.
    Only needs to exist long enough for add_picture() to embed its bytes —
    replaceImage swaps this out per-client at rewrite time."""
    from PIL import Image, ImageDraw

    img = Image.new("RGBA", (600, 200), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle([(0, 0), (599, 199)], radius=16, outline=(126, 200, 227, 255), width=6)
    draw.text((300, 100), "CLIENT LOGO", fill=(11, 31, 51, 255), anchor="mm")
    LOGO_PLACEHOLDER_PATH.parent.mkdir(parents=True, exist_ok=True)
    img.save(LOGO_PLACEHOLDER_PATH)
    return LOGO_PLACEHOLDER_PATH


def _set_alt_text(picture_shape, description: str) -> None:
    """Sets the picture's alt-text/description — this is what
    slides_rewriter.find_logo_placeholders() and Uncharted Ice's
    inspect_template.py both key off of ("logo"/"client_logo"/"client logo"
    substring match on title+description)."""
    picture_shape._element.nvPicPr.cNvPr.set("descr", description)
    picture_shape._element.nvPicPr.cNvPr.set("title", description)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank layout


def _fill_background(slide, color: RGBColor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, font=BODY_FONT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return box


def _bullets(slide, left, top, width, height, items, size, color, font=BODY_FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.name = font
            run.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def build_slide_1_cover(prs, logo_path: Path):
    slide = _blank_slide(prs)
    _fill_background(slide, NAVY)

    _textbox(slide, Inches(1), Inches(2.4), Inches(11.3), Inches(1.4),
              "CROSSING THE ICE\nINTERACTIVE KEYNOTE", 40, WHITE, bold=True, font=FONT)
    _textbox(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.6),
              "An interactive session by James Castrission", 20, ICE_BLUE, italic=True)
    _textbox(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5),
              "[EVENT_DATE]", 18, WHITE)
    _textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
              "inspire  ·  engage  ·  explore", 16, ICE_BLUE, align=PP_ALIGN.LEFT)

    pic = slide.shapes.add_picture(str(logo_path), Inches(10.4), Inches(0.5), height=Inches(0.9))
    _set_alt_text(pic, "client_logo")
    return slide


def build_slide_2_bio(prs):
    slide = _blank_slide(prs)
    _fill_background(slide, WHITE)
    _textbox(slide, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.7),
              "ONE OF AUSTRALIA'S MOST EXPERIENCED AND IN-DEMAND SPEAKERS",
              22, NAVY, bold=True, font=FONT)
    body = (
        "Hi, I'm James Castrission.\n\n"
        "I'm the holder of two Guinness World Records, I've written two best-selling "
        "books on adventuring, and in 2016, I was named one of Australia's 50 Greatest "
        "Explorers of All Time.\n\n"
        "I also spent five years as a management consultant at Deloitte and over the "
        "past decade, I have delivered keynotes and run workshops to over 980,000 "
        "people at 500+ corporate events in 40 different countries. That means, I "
        "understand what you are looking for from a keynote speaker at your next "
        "[CLIENT_ORG] event.\n\n"
        "Over the following pages, please find information showing how I can add "
        "value at your next conference.\n\n"
        "I look forward to working with you in creating an unforgettable conference."
    )
    _textbox(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), body, 15, CHARCOAL)
    return slide


def build_slide_3_positioning(prs):
    slide = _blank_slide(prs)
    _fill_background(slide, WHITE)
    _textbox(slide, Inches(0.8), Inches(0.6), Inches(8.2), Inches(1.0),
              "INSPIRE YOUR AUDIENCE:\nDREAM, DEFY, AND CONQUER", 24, NAVY, bold=True, font=FONT)
    body = (
        "Embark on an extraordinary journey with James Castrission, a world-record-"
        "breaking adventurer and internationally recognized motivational speaker, "
        "regarded as one of Australia's premier public speakers.\n\n"
        "Cas captivates audiences with inspiring stories that empower them to reach "
        "their full potential. As a record-breaking explorer and former Deloitte "
        "management consultant, he understands what it takes to succeed in both "
        "worlds—offering unique insights and proven strategies. Cas seamlessly "
        "blends adventure with motivation, delivering a truly unforgettable "
        "presentation."
    )
    _textbox(slide, Inches(0.8), Inches(1.8), Inches(8.0), Inches(5.0), body, 15, CHARCOAL)

    # Logistics sidebar
    sidebar_left, sidebar_top, sidebar_w, sidebar_h = Inches(9.2), Inches(1.8), Inches(3.3), Inches(4.6)
    sidebar = slide.shapes.add_shape(1, sidebar_left, sidebar_top, sidebar_w, sidebar_h)
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = NAVY
    sidebar.line.fill.background()
    tf = sidebar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    sidebar_lines = [
        ("DURATION", "[DURATION]"),
        ("DATE", "[EVENT_DATE]"),
        ("VENUE", "[VENUE]"),
        ("AUDIENCE SIZE", "[AUDIENCE_SIZE]"),
        ("SHOWREEL", "Click Play To Watch Showreel"),
    ]
    for i, (label, value) in enumerate(sidebar_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run1 = p.add_run()
        run1.text = f"{label}\n"
        run1.font.size = Pt(12)
        run1.font.bold = True
        run1.font.color.rgb = ICE_BLUE
        run1.font.name = FONT
        run2 = p.add_run()
        run2.text = value
        run2.font.size = Pt(14)
        run2.font.color.rgb = WHITE
        run2.font.name = BODY_FONT
    return slide


def build_slide_4_content(prs):
    slide = _blank_slide(prs)
    _fill_background(slide, WHITE)
    _textbox(slide, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6),
              "CROSSING THE ICE", 24, NAVY, bold=True, font=FONT)
    body = (
        "Crossing the Ice is a highly inspiring and stimulating talk, packed with "
        "practical insights on the power of collaboration and teamwork, mindset for "
        "success and navigating highly volatile environments.\n\n"
        "This keynote takes your team on a rich audio-visual journey through one of "
        "the greatest feats in modern exploration. In 100 years of polar exploration, "
        "no one had ever walked from the edge of Antarctica to the South Pole and back "
        "without any support. Many had tried. None had succeeded—not until Cas and "
        "his partner made history by completing the longest unsupported polar journey. "
        "They defied the odds, turning an audacious goal into reality through grit and "
        "teamwork.\n\n"
        "This high-energy and carefully crafted keynote will ensure your team is on "
        "the edge of their seats, empowered with takeaways and lessons they can apply "
        "to their professional and personal lives."
    )
    _textbox(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(4.3), body, 15, CHARCOAL)

    _textbox(slide, Inches(0.8), Inches(6.0), Inches(2.5), Inches(0.5),
              "LEARNING OUTCOMES", 14, NAVY, bold=True, font=FONT)
    _bullets(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.8),
             ["Resilience", "Working Together", "Adapting to Change", "Foundations of a High Performance Team"],
             13, CHARCOAL)
    return slide


def build_slide_5_testimonials(prs):
    slide = _blank_slide(prs)
    _fill_background(slide, WHITE)
    _textbox(slide, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6),
              "FEEDBACK FROM RECENT CLIENTS", 24, NAVY, bold=True, font=FONT)

    quotes = [
        ("“...Over the moon with James. Had the entire audience engaged... "
         "Everyone was raving about him.”", "Medical Meetings"),
        ("“...incredibly moving and effective exercise to rapidly help build an "
         "understanding of what's more important in team decision making under "
         "pressure.”", "Salesforce"),
        ("“...amazing immersive experience for our leadership team. very awe "
         "inspiring!”", "Uniting"),
    ]
    top = Inches(1.8)
    for quote, attribution in quotes:
        _textbox(slide, Inches(0.8), top, Inches(11.7), Inches(1.3), quote, 15, CHARCOAL, italic=True)
        _textbox(slide, Inches(0.8), top + Inches(1.0), Inches(11.7), Inches(0.4),
                  f"— {attribution}", 13, MUTED, bold=True)
        top += Inches(1.7)
    return slide


def build_slide_6_gifts(prs):
    slide = _blank_slide(prs)
    _fill_background(slide, WHITE)
    _textbox(slide, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6),
              "CONFERENCE GIFTS", 24, NAVY, bold=True, font=FONT)
    body = (
        "Aside from having an inspiring experience, many clients like their delegates "
        "to leave with a takeaway gift.\n\n"
        "These include adventure, branded conference clothing, my books and "
        "documentaries...\n\n"
        "The possibilities are endless! I'd love to work with your budget ensuring "
        "your delegates leave with a physical take away from the experience."
    )
    _textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.5), body, 16, CHARCOAL)
    return slide


def build_slide_7_pricing(prs):
    slide = _blank_slide(prs)
    _fill_background(slide, WHITE)
    _textbox(slide, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6),
              "YOUR INVESTMENT", 24, NAVY, bold=True, font=FONT)

    _textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4), "DESIGN", 16, NAVY, bold=True, font=FONT)
    _bullets(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.0), [
        "Customisation of an engaging 60 minute interactive keynote to maximise learning outcomes;",
        "Customised skill modules to match the desired programme outcomes delivered;",
        "Debrief calls with the programme team for feedback; and",
        "Prep calls for coordinating logistics.",
    ], 13, CHARCOAL)

    _textbox(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.4), "DELIVERY", 16, NAVY, bold=True, font=FONT)
    _bullets(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(3.0), [
        "Lead facilitator & adventurer James Castrission engage audience by linking "
        "together meeting facilitation to create a compelling impactful experience;",
        "Engaging and fun face-to-face delivery; and",
        "Tailored Insightful Thinking Questions.",
    ], 13, CHARCOAL)

    invest_box = slide.shapes.add_shape(1, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.4))
    invest_box.fill.solid()
    invest_box.fill.fore_color.rgb = NAVY
    invest_box.line.fill.background()
    tf = invest_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run1 = p.add_run()
    run1.text = "DESIGN & DELIVERY INVESTMENT\n"
    run1.font.size = Pt(13)
    run1.font.color.rgb = ICE_BLUE
    run1.font.bold = True
    run1.font.name = FONT
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "60 Minute Regular Investment - $8,500 AUD ex GST + Travel"
    run2.font.size = Pt(18)
    run2.font.bold = True
    run2.font.color.rgb = WHITE
    run2.font.name = FONT
    return slide


def build_slide_8_contact(prs):
    slide = _blank_slide(prs)
    _fill_background(slide, NAVY)
    _textbox(slide, Inches(1), Inches(2.4), Inches(11.3), Inches(1.0),
              "THANK YOU!", 40, WHITE, bold=True, font=FONT)
    contact = (
        "James Castrission I Lead Facilitator\n"
        "M: +61 402 904 334\n"
        "E: james@myadventuregroup.com.au\n"
        "W: www.myadventuregroup.com.au"
    )
    _textbox(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(2.0), contact, 16, ICE_BLUE)
    return slide


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    logo_path = _make_logo_placeholder_image()

    build_slide_1_cover(prs, logo_path)
    build_slide_2_bio(prs)
    build_slide_3_positioning(prs)
    build_slide_4_content(prs)
    build_slide_5_testimonials(prs)
    build_slide_6_gifts(prs)
    build_slide_7_pricing(prs)
    build_slide_8_contact(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


def upload(pptx_path: Path, name: str = "Interactive Keynote Template - Crossing the Ice") -> dict:
    """Uploads the built .pptx into the Templates Drive folder, converting it
    to a native Google Slides file (name must contain 'template' — see
    src/drive_service.get_master_template()). Requires a working .env
    (GOOGLE_CLIENT_ID/SECRET/REFRESH_TOKEN).

    get_master_template() expects exactly one matching file — if you're
    re-running this after a design tweak, delete/rename the previous
    template file in Drive first, or the pipeline will start failing with
    an "expected exactly one template" error."""
    from googleapiclient.http import MediaFileUpload

    import config as _config
    from src.google_clients import GoogleClients

    clients = GoogleClients()
    media = MediaFileUpload(
        str(pptx_path),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=False,
    )
    return (
        clients.drive.files()
        .create(
            body={
                "name": name,
                "parents": [_config.TEMPLATES_DRIVE_FOLDER_ID],
                "mimeType": "application/vnd.google-apps.presentation",
            },
            media_body=media,
            fields="id, name, mimeType, webViewLink",
            supportsAllDrives=True,
        )
        .execute()
    )


if __name__ == "__main__":
    path = build()
    print(f"Wrote master template: {path}")
    if "--upload" in sys.argv:
        result = upload(path)
        print("Uploaded and converted to Google Slides:", result)
